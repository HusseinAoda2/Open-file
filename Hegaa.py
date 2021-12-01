from tkinter import *
from tkinter import ttk
from tkinter import messagebox

# ----------- تثبيت ثم استدعاء مكتبة الإكسل -----------
# pip install openpyxl
from openpyxl import Workbook

# ----------- تثبيت ثم استدعاء مكتبة الوورد -----------
# pip install --pre python-docx
from docx import Document

# ----------- تثبيت ثم استدعاء مكتبة بوربوينت -----------
# pip install python-pptx
from pptx import Presentation

class Employee:
    # ----------- إنشاء نافذة البرنامج -----------
    def __init__(self, root):
        self.root = root
        self.root.geometry('300x350+610+260') # الأرقام الإضافية بعد الطول والعرض (22) هي توصيت الشاشة مسافة من اليسار ومن الأعلى
        self.root.title('إنشاء ملفات أوفيس')
        self.root.configure(background="#E2C2B9")
        self.root.resizable(False, False)
        title = Label(self.root,
                      text='[إنشاء ملف أوفيس]',
                      bg='#F2DDC1',
                      font=('monospace',18),
                      fg='#000'
                      )
        title.pack(fill='x')

        # ----------- تقسيم لأختيارات علوية -----------
        # nb = ttk.Notebook(self.root)
        # nb.place(x=900,y=34, width=400, height=460)

        # ----------- buttons الأزرار -----------
        # btn_Frame = Frame(self.root, bg="#fff")
        # btn_Frame.place(x=900, y=500, width=400, height=148)
        # title2 = Label(btn_Frame, text='لوحة التحكم', font=('monospace', 16), bg='#316df4', fg='#fff')
        # title2.pack(fill='x')

        excel_btn = Button(text='Excel File', bg='#99A799', fg='#D3E4CD', activebackground='#E2C2B9', cursor='plus', width=15, pady=10, padx=50, relief='ridge', command=self.excel)
        excel_btn.pack()

        word_btn = Button(text='Word File', bg='#99A799', fg='#D3E4CD', activebackground='#E2C2B9', width=15, pady=10, padx=50, relief='ridge', cursor='circle', command=self.word)
        word_btn.pack()

        power_btn = Button(text='PowerPoint File', bg='#99A799', fg='#D3E4CD', activebackground='#E2C2B9', width=15, pady=10, padx=50, relief='ridge', cursor='sizing', command=self.power)
        power_btn.pack()

        text_btn = Button(text='Text File', bg='#99A799', fg='#D3E4CD', activebackground='#E2C2B9', width=15, pady=10, padx=50, relief='ridge', cursor='target', command=self.text)
        text_btn.pack()

        exit_btn = Button(text='إغلاق البرنامج', bg='#99A799', fg='#D3E4CD', activebackground='#E2C2B9', width=15, pady=10, padx=50, relief='ridge', cursor='cross', command=root.quit)
        exit_btn.pack()

        about_btn = Button(text='من نحن', bg='#99A799', fg='#D3E4CD', activebackground='#E2C2B9', width=15, pady=10, padx=50, relief='ridge', cursor='heart', bitmap='info', command=self.about)
        about_btn.pack()

    # -------- دالة إنشاء ملف إكسل --------
    def excel(self):
        workbook = Workbook()
        sheet = workbook.active
        
        sheet["A2"] = "Email"
        sheet["B2"] = "Username"
        sheet["C2"] = "Password"
        sheet["D2"] = "Key"
        sheet["E2"] = "Skype"
        
        workbook.save(filename="Hegaa.xlsx")

    # -------- دالة إنشاء ملف وورد --------
    def word(self):
        document = Document()

        paragraph = document.add_paragraph('أهلا وسهلا بك \nتشرفت بالتواصل معك \n* الاسم : م. حسين محمد عودة\n* أعمل Team Leader لشركة Tech makers : \nhttps://techmakers.tech/\n - Web Devs : \n> Front End. \n> wordpress. \n- python developer \nأنشأت موقعي الشخصي وعليه كل ما يخصني من أعمال ومعلومات : \nhttps://husseinouda.me/\n* مهتم في البرمجة وخاصة لغة python.')

        document.save('Hegaa.docx')

    # -------- دالة إنشاء ملف بوربوينت --------
    def power(self):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Hello, Hegaa!"
        subtitle.text = "python-pptx was here!"

        prs.save('Hegaa.pptx')

    # -------- دالة إنشاء ملف نصي --------
    def text(self):
        file = open("Hegaa.txt", "w")

        file.write("\nأهلا وسهلا بك \nتشرفت بالتواصل معك \n* الاسم : م. حسين محمد عودة\n* أعمل Team Leader لشركة Tech makers : \nhttps://techmakers.tech/\n - Web Devs : \n> Front End. \n> wordpress. \n- python developer \nأنشأت موقعي الشخصي وعليه كل ما يخصني من أعمال ومعلومات : \nhttps://husseinouda.me/\n* مهتم في البرمجة وخاصة لغة python.")

        file.close()

    # -------- دالة من نحن --------
    def about(self):
        messagebox.showinfo("المطور حسين عودة", "husseinaoda@mail.com : مرحبا بكم في برنامج إنشاء ملفات أوفيس")


root = Tk()
ob = Employee(root)
root.mainloop()
